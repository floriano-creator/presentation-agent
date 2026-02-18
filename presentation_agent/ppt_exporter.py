"""PowerPoint export — high-end keynote-style design."""

from __future__ import annotations

import io
import tempfile
from pathlib import Path
from typing import Optional, Tuple

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from presentation_agent.design_library import (
    LayoutType,
    LINE_SPACING,
    get_template,
    select_layout,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    SLIDE_WIDTH_INCH,
    SLIDE_HEIGHT_INCH,
    ACCENT_BAR_WIDTH,
    ACCENT_LINE_HEIGHT,
    TITLE_AREA_MAX_HEIGHT,
    TITLE_BODY_GAP,
    MAX_BULLETS_DISPLAY,
    BODY_MIN_FONT_SIZE,
    BODY_FONT_REDUCE_IF_BULLETS_ABOVE,
)
from presentation_agent.models import ExportResult, SlideWithImage
from presentation_agent.theme_library import get_theme

# DPI assumption for image size (Pillow default)
_IMG_DPI = 96.0


class PPTExporterError(Exception):
    """Raised when PowerPoint export fails."""

    pass


def _get_image_size_inches(stream: io.BytesIO) -> Optional[Tuple[float, float]]:
    """Return (width_inch, height_inch) from image stream; None on failure."""
    try:
        from PIL import Image
        stream.seek(0)
        img = Image.open(stream)
        w_px, h_px = img.size
        stream.seek(0)
        if w_px <= 0 or h_px <= 0:
            return None
        return (w_px / _IMG_DPI, h_px / _IMG_DPI)
    except Exception:
        return None


class PPTExporter:
    """Exports presentation with theme-based styling."""

    def __init__(self) -> None:
        self._theme = None

    def _add_picture_fit_centered(
        self,
        slide,
        img_stream: io.BytesIO,
        left_inch: float,
        top_inch: float,
        box_width_inch: float,
        box_height_inch: float,
    ) -> bool:
        """
        Add image to slide preserving aspect ratio, fitted inside box and centered.
        Never stretches. Returns True if added, False on error.
        """
        size = _get_image_size_inches(img_stream)
        if not size:
            return False
        img_w, img_h = size
        if img_w <= 0 or img_h <= 0:
            return False
        scale = min(box_width_inch / img_w, box_height_inch / img_h)
        fit_w = img_w * scale
        fit_h = img_h * scale
        left_cen = left_inch + (box_width_inch - fit_w) / 2.0
        top_cen = top_inch + (box_height_inch - fit_h) / 2.0
        try:
            img_stream.seek(0)
            slide.shapes.add_picture(
                img_stream,
                Inches(left_cen),
                Inches(top_cen),
                width=Inches(fit_w),
                height=Inches(fit_h),
            )
            return True
        except Exception:
            return False

    def _apply_slide_background(self, slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self._theme.background)

    def _apply_accent_background(self, slide) -> None:
        fill = slide.background.fill
        if (
            self._theme.use_gradient
            and self._theme.gradient_start
            and self._theme.gradient_end
        ):
            try:
                fill.gradient()
                fill.gradient_angle = 135
                stops = fill.gradient_stops
                if len(stops) >= 1:
                    stops[0].color.rgb = RGBColor(*self._theme.gradient_start)
                    stops[0].position = 0
                if len(stops) >= 2:
                    stops[1].color.rgb = RGBColor(*self._theme.gradient_end)
                    stops[1].position = 1
            except Exception:
                fill.solid()
                fill.fore_color.rgb = RGBColor(*self._theme.section_bg)
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(*self._theme.section_bg)

    def _add_accent_bar_left(self, slide) -> None:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            Inches(ACCENT_BAR_WIDTH),
            SLIDE_HEIGHT,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*self._theme.accent)
        bar.line.fill.background()

    def _add_accent_line_top(self, slide) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            SLIDE_WIDTH,
            Inches(ACCENT_LINE_HEIGHT),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*self._theme.accent)
        line.line.fill.background()

    def _add_overlay(self, slide, alpha: float = 0.6) -> None:
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            SLIDE_WIDTH,
            SLIDE_HEIGHT,
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(*self._theme.overlay_dark)
        overlay.fill.transparency = alpha
        overlay.line.fill.background()

    def _add_card_panel(self, slide, left: float, top: float, width: float, height: float) -> None:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*self._theme.card_bg)
        card.line.fill.background()
        if hasattr(card, "adjustments") and card.adjustments:
            card.adjustments[0] = 0.05

    def _apply_title_style(self, paragraph, size: Optional[int] = None) -> None:
        paragraph.font.size = Pt(size or self._theme.title_size)
        paragraph.font.bold = True
        paragraph.font.name = self._theme.font_family
        paragraph.font.color.rgb = RGBColor(*self._theme.primary)

    def _apply_body_style(self, paragraph) -> None:
        paragraph.font.size = Pt(self._theme.body_size)
        paragraph.font.name = self._theme.font_family
        paragraph.font.color.rgb = RGBColor(*self._theme.primary)
        paragraph.space_after = Pt(12)
        paragraph.line_spacing = LINE_SPACING

    def _apply_caption_style(self, paragraph) -> None:
        paragraph.font.size = Pt(self._theme.caption_size)
        paragraph.font.name = self._theme.font_family
        paragraph.font.color.rgb = RGBColor(*self._theme.secondary)

    def _apply_section_text(self, paragraph) -> None:
        paragraph.font.name = self._theme.font_family
        paragraph.font.color.rgb = RGBColor(*self._theme.section_text)

    def _prepare_title(self, title: str, max_chars: int = 120) -> str:
        """Keep title in safe zone; truncate with ellipsis if too long."""
        if not title:
            return title
        title = title.strip()
        if len(title) <= max_chars:
            return title
        return title[: max_chars - 1].rstrip() + "…"

    def _prepare_bullets(
        self, bullet_points: list[str]
    ) -> Tuple[list[str], Optional[int]]:
        """
        Cap bullets for layout safety. Returns (list to render, optional font size override).
        """
        if not bullet_points:
            return [], None
        bullets = bullet_points[:MAX_BULLETS_DISPLAY]
        font_override = None
        if len(bullet_points) >= BODY_FONT_REDUCE_IF_BULLETS_ABOVE:
            body_size = getattr(self._theme, "body_size", 18)
            font_override = max(BODY_MIN_FONT_SIZE, body_size - 2)
        return bullets, font_override

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = "") -> None:
        template = get_template(LayoutType.TITLE_SLIDE)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._apply_slide_background(slide)

        if template.accent_top:
            self._add_accent_line_top(slide)

        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(title)
            p.alignment = 1
            self._apply_title_style(p, size=52)

        if template.subtitle_area and subtitle:
            l, t, w, h = template.subtitle_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.alignment = 1
            self._apply_caption_style(p)

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: SlideWithImage,
        slide_index: int,
        total_slides: int,
    ) -> int:
        layout = select_layout(slide_data, slide_index, total_slides)
        template = get_template(layout)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        images_added = 0

        if layout == LayoutType.BOLD_SECTION_DIVIDER:
            self._apply_accent_background(slide)
            self._render_bold_section(slide, template, slide_data)
        elif layout == LayoutType.HERO_BACKGROUND and slide_data.image_url:
            images_added = self._render_hero_background(slide, template, slide_data)
        elif layout == LayoutType.HERO_RIGHT:
            images_added = self._render_hero_right(slide, template, slide_data)
        elif layout == LayoutType.CARD_LAYOUT:
            self._apply_slide_background(slide)
            self._render_card_layout(slide, template, slide_data)
        elif layout == LayoutType.ACCENT_LEFT_LAYOUT or layout == LayoutType.CONCLUSION:
            self._apply_slide_background(slide)
            self._add_accent_bar_left(slide)
            self._render_title_and_bullets(slide, template, slide_data)
        elif layout == LayoutType.MINIMAL_TEXT:
            self._apply_slide_background(slide)
            self._render_minimal_text(slide, template, slide_data)
        elif layout == LayoutType.IMAGE_FOCUS:
            self._apply_slide_background(slide)
            images_added = self._render_image_focus(slide, template, slide_data)
        elif layout == LayoutType.TWO_COLUMN:
            self._apply_slide_background(slide)
            images_added = self._render_two_column(slide, template, slide_data)
        else:
            self._apply_slide_background(slide)
            if template.accent_left:
                self._add_accent_bar_left(slide)
            self._render_title_and_bullets(slide, template, slide_data)

        if slide_data.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data.speaker_notes

        return images_added

    def _render_bold_section(
        self, slide, template, slide_data: SlideWithImage
    ) -> None:
        """BOLD_SECTION_DIVIDER — large centered title, accent bg."""
        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            p.alignment = 1
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.name = self._theme.font_family
            self._apply_section_text(p)

        bullets, _ = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bullets[0]
            p.alignment = 1
            p.font.size = Pt(self._theme.body_size)
            p.font.name = self._theme.font_family
            self._apply_section_text(p)

    def _render_hero_background(
        self, slide, template, slide_data: SlideWithImage
    ) -> int:
        """HERO_BACKGROUND — full image (aspect ratio preserved), overlay, text on top."""
        try:
            img_stream = self._fetch_image(slide_data.image_url)
            if img_stream and not self._add_picture_fit_centered(
                slide, img_stream, 0, 0, SLIDE_WIDTH_INCH, SLIDE_HEIGHT_INCH
            ):
                img_stream = None
            if not img_stream:
                self._apply_slide_background(slide)
                return 0
        except Exception:
            self._apply_slide_background(slide)
            return 0

        self._add_overlay(slide, alpha=0.55)

        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            p.alignment = 1
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.name = self._theme.font_family
            self._apply_section_text(p)

        bullets, body_font = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            sz = body_font or self._theme.body_size
            for i, point in enumerate(bullets[:3]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                p.font.size = Pt(sz)
                p.font.name = self._theme.font_family
                self._apply_section_text(p)
                p.space_after = Pt(8)

        return 1

    def _render_hero_right(
        self, slide, template, slide_data: SlideWithImage
    ) -> int:
        """HERO_RIGHT — text left, full-height image right."""
        images_added = 0
        self._apply_slide_background(slide)

        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            self._apply_title_style(p)

        bullets, body_font = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                self._apply_body_style(p)
                if body_font:
                    p.font.size = Pt(body_font)
                    p.space_after = Pt(4)

        if template.image_area and slide_data.image_url:
            try:
                img_stream = self._fetch_image(slide_data.image_url)
                if img_stream:
                    b = template.image_area
                    if self._add_picture_fit_centered(
                        slide, img_stream, b.left, b.top, b.width, b.height
                    ):
                        images_added = 1
            except Exception:
                pass

        return images_added

    def _render_card_layout(
        self, slide, template, slide_data: SlideWithImage
    ) -> None:
        """CARD_LAYOUT — text inside light rounded panel."""
        if template.body_area:
            b = template.body_area
            self._add_card_panel(slide, b.left - 0.2, b.top - 0.2, b.width + 0.4, b.height + 0.4)

        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            self._apply_title_style(p)

        bullets, body_font = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                self._apply_body_style(p)
                if body_font:
                    p.font.size = Pt(body_font)
                    p.space_after = Pt(4)

    def _render_minimal_text(
        self, slide, template, slide_data: SlideWithImage
    ) -> None:
        """MINIMAL_TEXT — very large title, minimal body."""
        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            p.alignment = 1
            self._apply_title_style(p, size=44)

        bullets, _ = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                p.alignment = 1
                self._apply_body_style(p)

    def _render_title_and_bullets(
        self, slide, template, slide_data: SlideWithImage
    ) -> None:
        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            self._apply_title_style(p)

        bullets, body_font = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                self._apply_body_style(p)
                if body_font:
                    p.font.size = Pt(body_font)
                    p.space_after = Pt(4)

    def _render_two_column(
        self, slide, template, slide_data: SlideWithImage
    ) -> int:
        images_added = 0
        if template.title_area:
            l, t, w, h = template.title_area.inches
            box = slide.shapes.add_textbox(l, t, w, min(h, Inches(TITLE_AREA_MAX_HEIGHT)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._prepare_title(slide_data.title)
            self._apply_title_style(p)

        bullets, body_font = self._prepare_bullets(slide_data.bullet_points or [])
        if template.body_area and bullets:
            l, t, w, h = template.body_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"•  {point}"
                self._apply_body_style(p)
                if body_font:
                    p.font.size = Pt(body_font)
                    p.space_after = Pt(4)

        if template.image_area and slide_data.image_url:
            try:
                img_stream = self._fetch_image(slide_data.image_url)
                if img_stream:
                    b = template.image_area
                    if self._add_picture_fit_centered(
                        slide, img_stream, b.left, b.top, b.width, b.height
                    ):
                        images_added = 1
            except Exception:
                pass
        return images_added

    def _render_image_focus(
        self, slide, template, slide_data: SlideWithImage
    ) -> int:
        images_added = 0
        if template.image_area and slide_data.image_url:
            try:
                img_stream = self._fetch_image(slide_data.image_url)
                if img_stream:
                    b = template.image_area
                    if self._add_picture_fit_centered(
                        slide, img_stream, b.left, b.top, b.width, b.height
                    ):
                        images_added = 1
            except Exception:
                pass

        if template.caption_area:
            l, t, w, h = template.caption_area.inches
            box = slide.shapes.add_textbox(l, t, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            cap = self._prepare_title(slide_data.title, max_chars=80)
            if slide_data.bullet_points:
                cap += " — " + (slide_data.bullet_points[0][:60] + "…" if len(slide_data.bullet_points[0]) > 60 else slide_data.bullet_points[0])
            p.text = cap
            self._apply_caption_style(p)
        return images_added

    def _fetch_image(self, url: str) -> Optional[io.BytesIO]:
        with httpx.Client(timeout=10.0) as client:
            response = client.get(url)
            response.raise_for_status()
            return io.BytesIO(response.content)

    def export(
        self,
        slides: list[SlideWithImage],
        title: str,
        output_path: Optional[Path | str] = None,
        theme: Optional[str] = None,
    ) -> ExportResult:
        self._theme = get_theme(theme)

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        self._add_title_slide(prs, title)

        images_included = 0
        total = len(slides)
        for i, s in enumerate(slides):
            images_included += self._add_content_slide(
                prs, s, slide_index=i, total_slides=total
            )

        if output_path is None:
            fd, path = tempfile.mkstemp(suffix=".pptx")
            import os
            os.close(fd)
            output_path = path

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))

        return ExportResult(
            output_path=str(path.resolve()),
            slide_count=len(slides) + 1,
            images_included=images_included,
        )
